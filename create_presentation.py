from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_ageless_presentation():
    """株式会社エイジレスの会社紹介プレゼンテーションを作成"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # スライド1: タイトル
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
    
    # 背景グラデーション（紫色系）
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(102, 126, 234)
    background.line.fill.background()
    
    # タイトル
    title_box = slide1.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "株式会社エイジレス"
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    subtitle_box = slide1.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Ageless Inc. - 企業紹介"
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # スライド2: 会社概要
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background2.fill.solid()
    background2.fill.fore_color.rgb = RGBColor(245, 87, 108)
    background2.line.fill.background()
    
    # 見出し
    heading2 = slide2.shapes.add_textbox(Inches(0), Inches(0.8), prs.slide_width, Inches(1))
    heading2_frame = heading2.text_frame
    heading2_frame.text = "会社概要"
    heading2_frame.paragraphs[0].font.size = Pt(48)
    heading2_frame.paragraphs[0].font.bold = True
    heading2_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    heading2_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 内容
    content2 = slide2.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(4))
    content2_frame = content2.text_frame
    content2_frame.word_wrap = True
    
    about_text = [
        "株式会社エイジレスは、年齢にとらわれない価値の創造を使命とする企業です。",
        "",
        "テクノロジーと人間性の融合により、",
        "すべての世代が活躍できる社会の実現を目指しています。",
        "",
        "設立: 20XX年",
        "従業員数: XX名",
        "資本金: X,XXX万円"
    ]
    
    for i, text in enumerate(about_text):
        if i == 0:
            p = content2_frame.paragraphs[0]
        else:
            p = content2_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(24) if text else Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)
    
    # スライド3: ビジョン
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background3.fill.solid()
    background3.fill.fore_color.rgb = RGBColor(79, 172, 254)
    background3.line.fill.background()
    
    heading3 = slide3.shapes.add_textbox(Inches(0), Inches(0.8), prs.slide_width, Inches(1))
    heading3_frame = heading3.text_frame
    heading3_frame.text = "ビジョン"
    heading3_frame.paragraphs[0].font.size = Pt(48)
    heading3_frame.paragraphs[0].font.bold = True
    heading3_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    heading3_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ビジョンボックス
    vision_box = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.5), Inches(9.333), Inches(3)
    )
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    vision_box.fill.fore_color.brightness = 0.2
    vision_box.line.fill.background()
    
    vision_text = slide3.shapes.add_textbox(Inches(2.5), Inches(3), Inches(8.333), Inches(2.5))
    vision_frame = vision_text.text_frame
    vision_frame.word_wrap = True
    
    vision_lines = [
        "年齢という境界を超え、",
        "すべての人が可能性を最大限に発揮できる",
        "世界の創造"
    ]
    
    for i, line in enumerate(vision_lines):
        if i == 0:
            p = vision_frame.paragraphs[0]
        else:
            p = vision_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)
    
    # スライド4: 事業内容
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background4.fill.solid()
    background4.fill.fore_color.rgb = RGBColor(67, 233, 123)
    background4.line.fill.background()
    
    heading4 = slide4.shapes.add_textbox(Inches(0), Inches(0.5), prs.slide_width, Inches(1))
    heading4_frame = heading4.text_frame
    heading4_frame.text = "事業内容"
    heading4_frame.paragraphs[0].font.size = Pt(48)
    heading4_frame.paragraphs[0].font.bold = True
    heading4_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    heading4_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 3つのサービスカード
    services = [
        ("🚀 デジタル変革", "最新技術を活用した\nビジネス変革支援"),
        ("👥 人材開発", "多世代共生型の\n人材育成プログラム"),
        ("🌐 グローバル展開", "国境を越えた\nビジネスサポート")
    ]
    
    card_width = Inches(3.5)
    card_height = Inches(3.5)
    start_x = Inches(1.2)
    gap = Inches(0.5)
    
    for i, (title, desc) in enumerate(services):
        x = start_x + i * (card_width + gap)
        
        # カード背景
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.fill.background()
        
        # タイトル
        title_box = slide4.shapes.add_textbox(x, Inches(2.8), card_width, Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 説明
        desc_box = slide4.shapes.add_textbox(x + Inches(0.2), Inches(3.8), card_width - Inches(0.4), Inches(1.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(18)
        desc_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # スライド5: 強み
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background5.fill.solid()
    background5.fill.fore_color.rgb = RGBColor(250, 112, 154)
    background5.line.fill.background()
    
    heading5 = slide5.shapes.add_textbox(Inches(0), Inches(0.5), prs.slide_width, Inches(1))
    heading5_frame = heading5.text_frame
    heading5_frame.text = "私たちの強み"
    heading5_frame.paragraphs[0].font.size = Pt(48)
    heading5_frame.paragraphs[0].font.bold = True
    heading5_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    heading5_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    strengths = [
        "多様な世代の知見を活かしたイノベーション",
        "最新テクノロジーと実務経験の融合",
        "お客様に寄り添う柔軟な対応力",
        "持続可能な社会への貢献"
    ]
    
    for i, strength in enumerate(strengths):
        y = Inches(2 + i * 1.2)
        
        # 番号の丸
        circle = slide5.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(2), y, Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(26, 26, 46)
        circle.line.fill.background()
        
        # 番号
        num_box = slide5.shapes.add_textbox(Inches(2), y, Inches(0.8), Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = str(i + 1)
        num_frame.paragraphs[0].font.size = Pt(28)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # テキスト
        text_box = slide5.shapes.add_textbox(Inches(3.2), y + Inches(0.15), Inches(8), Inches(0.8))
        text_frame = text_box.text_frame
        text_frame.text = strength
        text_frame.paragraphs[0].font.size = Pt(24)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    
    # スライド6: お問い合わせ
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background6 = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background6.fill.solid()
    background6.fill.fore_color.rgb = RGBColor(168, 237, 234)
    background6.line.fill.background()
    
    heading6 = slide6.shapes.add_textbox(Inches(0), Inches(0.8), prs.slide_width, Inches(1))
    heading6_frame = heading6.text_frame
    heading6_frame.text = "お問い合わせ"
    heading6_frame.paragraphs[0].font.size = Pt(48)
    heading6_frame.paragraphs[0].font.bold = True
    heading6_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    heading6_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    contact_info = [
        ("会社名", "株式会社エイジレス"),
        ("所在地", "〒XXX-XXXX 東京都XXX区XXX"),
        ("TEL", "03-XXXX-XXXX"),
        ("Email", "info@ageless.co.jp"),
        ("Web", "www.ageless.co.jp")
    ]
    
    for i, (label, value) in enumerate(contact_info):
        y = Inches(2.2 + i * 0.9)
        
        # ラベル
        label_box = slide6.shapes.add_textbox(Inches(3.5), y, Inches(2.5), Inches(0.6))
        label_frame = label_box.text_frame
        label_frame.text = label + ":"
        label_frame.paragraphs[0].font.size = Pt(22)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # 値
        value_box = slide6.shapes.add_textbox(Inches(6.2), y, Inches(5), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(22)
        value_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    
    # 保存
    prs.save('/root/.openclaw/workspace/language_rpg/株式会社エイジレス_会社紹介.pptx')
    print("プレゼンテーションを作成しました: 株式会社エイジレス_会社紹介.pptx")

if __name__ == "__main__":
    create_ageless_presentation()
