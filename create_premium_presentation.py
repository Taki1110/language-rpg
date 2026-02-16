from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from lxml import etree
import io

def add_gradient_background(slide, color1, color2, direction="linear"):
    """グラデーション背景を追加"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    
    # ソリッドフィルで代用（python-pptxの制限）
    background.fill.solid()
    background.fill.fore_color.rgb = color1
    background.line.fill.background()
    
    # 2つ目の色でオーバーレイを作成して擬似的なグラデーション
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = color2
    # 透明度を設定（50%）
    spPr = overlay._element.spPr
    fill = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if fill is not None:
        srgbClr = fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '50000')  # 50%透明度
    overlay.line.fill.background()

def add_shadow_textbox(slide, left, top, width, height, text, font_size, font_color, 
                       shadow_color=RGBColor(0, 0, 0), shadow_offset=3):
    """シャドウ付きテキストボックス"""
    # シャドウ（オフセットして先に描画）
    shadow = slide.shapes.add_textbox(
        left + Pt(shadow_offset), top + Pt(shadow_offset), width, height
    )
    shadow_frame = shadow.text_frame
    shadow_frame.text = text
    shadow_frame.paragraphs[0].font.size = font_size
    shadow_frame.paragraphs[0].font.bold = True
    shadow_frame.paragraphs[0].font.color.rgb = shadow_color
    shadow_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # メインテキスト
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = font_size
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = font_color
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return textbox

def add_icon_circle(slide, left, top, size, number, bg_color, text_color):
    """アイコン付き丸"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.fill.background()
    
    # 番号テキスト
    textbox = slide.shapes.add_textbox(left, top + Pt(8), size, size)
    text_frame = textbox.text_frame
    text_frame.text = str(number)
    text_frame.paragraphs[0].font.size = Pt(36)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = text_color
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return circle

def add_3d_card(slide, left, top, width, height, title, description, 
                title_color=RGBColor(102, 126, 234), accent_color=RGBColor(255, 200, 100)):
    """3D風カード"""
    # シャドウ（擬似3D効果）
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Pt(4), top + Pt(4), width, height
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow.line.fill.background()
    spPr = shadow._element.spPr
    fill = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if fill is not None:
        srgbClr = fill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '30000')  # 30%透明度
    
    # メインカード
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(220, 220, 220)
    card.line.width = Pt(1)
    
    # アクセントバー（上部）
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(8)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        left + Pt(15), top + Pt(25), width - Pt(30), Pt(40)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = title_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 説明
    desc_box = slide.shapes.add_textbox(
        left + Pt(15), top + Pt(70), width - Pt(30), height - Pt(90)
    )
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
    desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    desc_frame.paragraphs[0].line_spacing = 1.5

def add_divider_line(slide, left, top, width, color, thickness=3):
    """装飾用ライン"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def create_premium_presentation():
    """プレミアム品質のプレゼンテーションを作成"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # カラーパレット
    PRIMARY_PURPLE = RGBColor(102, 126, 234)
    SECONDARY_PINK = RGBColor(245, 87, 108)
    ACCENT_GOLD = RGBColor(255, 200, 100)
    DARK_BG = RGBColor(30, 30, 46)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 245)
    
    # ========== スライド1: タイトル ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # グラデーション背景
    add_gradient_background(slide1, PRIMARY_PURPLE, SECONDARY_PINK)
    
    # 装飾ライン
    add_divider_line(slide1, Inches(4), Inches(2.2), Inches(5.333), ACCENT_GOLD, 4)
    
    # メインタイトル（シャドウ付き）
    add_shadow_textbox(
        slide1, Inches(0), Inches(2.5), Inches(13.333), Inches(1.2),
        "株式会社エイジレス", Pt(66), WHITE, RGBColor(0, 0, 0), 4
    )
    
    # サブタイトル
    subtitle = slide1.shapes.add_textbox(
        Inches(0), Inches(4), Inches(13.333), Inches(0.8)
    )
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "Ageless Inc."
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_GOLD
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.paragraphs[0].font.italic = True
    
    # キャッチコピー
    catch = slide1.shapes.add_textbox(
        Inches(0), Inches(5), Inches(13.333), Inches(0.6)
    )
    catch_frame = catch.text_frame
    catch_frame.text = "年齢を超えた価値創造"
    catch_frame.paragraphs[0].font.size = Pt(20)
    catch_frame.paragraphs[0].font.color.rgb = WHITE
    catch_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    add_divider_line(slide1, Inches(4), Inches(5.8), Inches(5.333), ACCENT_GOLD, 4)
    
    # ========== スライド2: 会社概要 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_GRAY
    bg2.line.fill.background()
    
    # ヘッダーバー
    header2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    header2.fill.solid()
    header2.fill.fore_color.rgb = PRIMARY_PURPLE
    header2.line.fill.background()
    
    # タイトル
    title2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    title2_frame = title2.text_frame
    title2_frame.text = "会社概要"
    title2_frame.paragraphs[0].font.size = Pt(40)
    title2_frame.paragraphs[0].font.bold = True
    title2_frame.paragraphs[0].font.color.rgb = WHITE
    
    # メインコンテンツエリア
    content_bg = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.8), Inches(11.333), Inches(5.2)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = WHITE
    content_bg.line.color.rgb = RGBColor(220, 220, 220)
    
    # ミッションステートメント
    mission_title = slide2.shapes.add_textbox(
        Inches(1.5), Inches(2.2), Inches(10), Inches(0.6)
    )
    mt_frame = mission_title.text_frame
    mt_frame.text = "ミッション"
    mt_frame.paragraphs[0].font.size = Pt(24)
    mt_frame.paragraphs[0].font.bold = True
    mt_frame.paragraphs[0].font.color.rgb = PRIMARY_PURPLE
    
    mission_text = slide2.shapes.add_textbox(
        Inches(1.5), Inches(2.8), Inches(10), Inches(1.5)
    )
    mission_frame = mission_text.text_frame
    mission_frame.word_wrap = True
    mission_content = [
        "年齢にとらわれない価値の創造を使命とし、",
        "テクノロジーと人間性の融合により、",
        "すべての世代が活躍できる社会の実現を目指す"
    ]
    for i, line in enumerate(mission_content):
        if i == 0:
            p = mission_frame.paragraphs[0]
        else:
            p = mission_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(60, 60, 60)
        p.space_before = Pt(6)
    
    # 会社情報テーブル風
    info_data = [
        ("設立", "20XX年"),
        ("従業員数", "XX名"),
        ("資本金", "X,XXX万円"),
        ("本社所在地", "東京都XXX区XXX")
    ]
    
    y_start = Inches(4.5)
    for i, (label, value) in enumerate(info_data):
        y = y_start + i * Inches(0.5)
        
        # ラベル
        label_box = slide2.shapes.add_textbox(Inches(2), y, Inches(3), Inches(0.4))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(16)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = PRIMARY_PURPLE
        
        # セパレータ
        sep = slide2.shapes.add_textbox(Inches(4.5), y, Inches(0.5), Inches(0.4))
        sep_frame = sep.text_frame
        sep_frame.text = ":"
        sep_frame.paragraphs[0].font.size = Pt(16)
        sep_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        
        # 値
        value_box = slide2.shapes.add_textbox(Inches(5), y, Inches(6), Inches(0.4))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(16)
        value_frame.paragraphs[0].font.color.rgb = RGBColor(60, 60, 60)
    
    # ========== スライド3: ビジョン ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_gradient_background(slide3, RGBColor(79, 172, 254), RGBColor(0, 242, 254))
    
    # ヘッダー
    header3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header3.fill.solid()
    header3.fill.fore_color.rgb = RGBColor(0, 0, 0)
    spPr3 = header3._element.spPr
    fill3 = spPr3.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if fill3 is not None:
        srgbClr3 = fill3.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr3 is not None:
            alpha3 = etree.SubElement(srgbClr3, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha3.set('val', '40000')  # 40%透明度
    header3.line.fill.background()
    
    # タイトル
    title3 = slide3.shapes.add_textbox(
        Inches(0), Inches(0.35), Inches(13.333), Inches(0.7)
    )
    title3_frame = title3.text_frame
    title3_frame.text = "ビジョン"
    title3_frame.paragraphs[0].font.size = Pt(40)
    title3_frame.paragraphs[0].font.bold = True
    title3_frame.paragraphs[0].font.color.rgb = WHITE
    title3_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ビジョンカード（中央配置）
    vision_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.2), Inches(8.333), Inches(3.5)
    )
    vision_card.fill.solid()
    vision_card.fill.fore_color.rgb = WHITE
    vision_card.line.fill.background()
    
    # ビジョンアイコン
    icon_bg = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5)
    )
    icon_bg.fill.solid()
    icon_bg.fill.fore_color.rgb = ACCENT_GOLD
    icon_bg.line.fill.background()
    
    icon_text = slide3.shapes.add_textbox(Inches(5.9), Inches(1.8), Inches(1.5), Inches(1))
    icon_tf = icon_text.text_frame
    icon_tf.text = "★"
    icon_tf.paragraphs[0].font.size = Pt(48)
    icon_tf.paragraphs[0].font.color.rgb = WHITE
    icon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ビジョンテキスト
    vision_lines = [
        "年齢という境界を超え、",
        "すべての人が可能性を",
        "最大限に発揮できる",
        "世界の創造"
    ]
    
    vy = Inches(2.8)
    for i, line in enumerate(vision_lines):
        vbox = slide3.shapes.add_textbox(Inches(3), vy + i * Inches(0.5), Inches(7.333), Inches(0.5))
        vf = vbox.text_frame
        vf.text = line
        vf.paragraphs[0].font.size = Pt(26) if i == 3 else Pt(22)
        vf.paragraphs[0].font.bold = True if i == 3 else False
        vf.paragraphs[0].font.color.rgb = PRIMARY_PURPLE if i == 3 else RGBColor(60, 60, 60)
        vf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== スライド4: 事業内容 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = RGBColor(245, 247, 250)
    bg4.line.fill.background()
    
    # ヘッダー
    header4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    header4.fill.solid()
    header4.fill.fore_color.rgb = RGBColor(67, 233, 123)
    header4.line.fill.background()
    
    title4 = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    title4_frame = title4.text_frame
    title4_frame.text = "事業内容"
    title4_frame.paragraphs[0].font.size = Pt(40)
    title4_frame.paragraphs[0].font.bold = True
    title4_frame.paragraphs[0].font.color.rgb = RGBColor(26, 46, 26)
    
    # 3つのサービスカード
    services = [
        {
            "icon": "🚀",
            "title": "デジタル変革",
            "desc": "AI・クラウド・データ分析など\n最新技術を活用し、\n企業のビジネスプロセスを\n根本から変革します",
            "accent": RGBColor(102, 126, 234)
        },
        {
            "icon": "👥",
            "title": "人材開発",
            "desc": "多世代が学び合う\n共生型の人材育成プログラムで、\n組織の持続的成長を\n支援します",
            "accent": RGBColor(245, 87, 108)
        },
        {
            "icon": "🌐",
            "title": "グローバル展開",
            "desc": "言語と文化の壁を超えた\nグローバルビジネス支援で、\n世界市場での成功を\n目指します",
            "accent": RGBColor(255, 200, 100)
        }
    ]
    
    card_width = Inches(3.8)
    card_height = Inches(4.5)
    start_x = Inches(0.8)
    gap = Inches(0.4)
    
    for i, service in enumerate(services):
        x = start_x + i * (card_width + gap)
        
        # カード背景（シャドウ）
        shadow = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Pt(6), Inches(1.9) + Pt(6), card_width, card_height
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
        shadow.line.fill.background()
        spPr_s = shadow._element.spPr
        fill_s = spPr_s.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if fill_s is not None:
            srgbClr_s = fill_s.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if srgbClr_s is not None:
                alpha_s = etree.SubElement(srgbClr_s, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                alpha_s.set('val', '15000')
        
        # メインカード
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()
        
        # アクセントバー
        accent = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, Inches(1.9), card_width, Pt(10)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = service["accent"]
        accent.line.fill.background()
        
        # アイコン
        icon_box = slide4.shapes.add_textbox(x, Inches(2.4), card_width, Inches(0.8))
        icon_frame = icon_box.text_frame
        icon_frame.text = service["icon"]
        icon_frame.paragraphs[0].font.size = Pt(48)
        icon_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # タイトル
        st_box = slide4.shapes.add_textbox(x, Inches(3.3), card_width, Inches(0.6))
        st_frame = st_box.text_frame
        st_frame.text = service["title"]
        st_frame.paragraphs[0].font.size = Pt(26)
        st_frame.paragraphs[0].font.bold = True
        st_frame.paragraphs[0].font.color.rgb = service["accent"]
        st_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 説明
        sd_box = slide4.shapes.add_textbox(
            x + Pt(15), Inches(4), card_width - Pt(30), Inches(2)
        )
        sd_frame = sd_box.text_frame
        sd_frame.word_wrap = True
        sd_frame.text = service["desc"]
        sd_frame.paragraphs[0].font.size = Pt(14)
        sd_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
        sd_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sd_frame.paragraphs[0].line_spacing = 1.4
    
    # ========== スライド5: 強み ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_gradient_background(slide5, RGBColor(250, 112, 154), RGBColor(254, 225, 64))
    
    # ヘッダー
    header5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    header5.fill.solid()
    header5.fill.fore_color.rgb = RGBColor(26, 26, 46)
    header5.line.fill.background()
    
    title5 = slide5.shapes.add_textbox(
        Inches(0), Inches(0.3), Inches(13.333), Inches(0.8)
    )
    title5_frame = title5.text_frame
    title5_frame.text = "私たちの強み"
    title5_frame.paragraphs[0].font.size = Pt(40)
    title5_frame.paragraphs[0].font.bold = True
    title5_frame.paragraphs[0].font.color.rgb = WHITE
    title5_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    strengths = [
        ("01", "多様な世代の知見を\n活かしたイノベーション", PRIMARY_PURPLE),
        ("02", "最新テクノロジーと\n実務経験の融合", SECONDARY_PINK),
        ("03", "お客様に寄り添う\n柔軟な対応力", ACCENT_GOLD),
        ("04", "持続可能な社会への\n貢献", RGBColor(67, 233, 123))
    ]
    
    # 2x2グリッド配置
    positions = [
        (Inches(1.5), Inches(1.8)),
        (Inches(7), Inches(1.8)),
        (Inches(1.5), Inches(4.3)),
        (Inches(7), Inches(4.3))
    ]
    
    for i, (num, text, color) in enumerate(strengths):
        x, y = positions[i]
        
        # 番号サークル
        circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.2), Inches(1.2))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        
        num_box = slide5.shapes.add_textbox(x, y + Pt(12), Inches(1.2), Inches(0.8))
        num_frame = num_box.text_frame
        num_frame.text = num
        num_frame.paragraphs[0].font.size = Pt(28)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = WHITE
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # テキスト
        text_box = slide5.shapes.add_textbox(x + Inches(1.4), y + Pt(10), Inches(4), Inches(1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 26, 46)
    
    # ========== スライド6: お問い合わせ ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg6 = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = RGBColor(168, 237, 234)
    bg6.line.fill.background()
    
    # 装飾パターン（上部）
    pattern = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.5)
    )
    pattern.fill.solid()
    pattern.fill.fore_color.rgb = PRIMARY_PURPLE
    pattern.line.fill.background()
    
    # タイトル
    title6 = slide6.shapes.add_textbox(
        Inches(0), Inches(0.6), Inches(13.333), Inches(1)
    )
    title6_frame = title6.text_frame
    title6_frame.text = "お問い合わせ"
    title6_frame.paragraphs[0].font.size = Pt(48)
    title6_frame.paragraphs[0].font.bold = True
    title6_frame.paragraphs[0].font.color.rgb = WHITE
    title6_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 連絡先カード
    contact_card = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2.8), Inches(7.333), Inches(4)
    )
    contact_card.fill.solid()
    contact_card.fill.fore_color.rgb = WHITE
    contact_card.line.fill.background()
    
    contact_info = [
        ("会社名", "株式会社エイジレス", PRIMARY_PURPLE),
        ("所在地", "〒XXX-XXXX 東京都XXX区XXX", RGBColor(60, 60, 60)),
        ("TEL", "03-XXXX-XXXX", RGBColor(60, 60, 60)),
        ("Email", "info@ageless.co.jp", RGBColor(60, 60, 60)),
        ("Web", "www.ageless.co.jp", PRIMARY_PURPLE)
    ]
    
    cy = Inches(3.2)
    for label, value, vcolor in contact_info:
        # ラベル
        cl_box = slide6.shapes.add_textbox(Inches(3.5), cy, Inches(2.5), Inches(0.5))
        cl_frame = cl_box.text_frame
        cl_frame.text = label
        cl_frame.paragraphs[0].font.size = Pt(16)
        cl_frame.paragraphs[0].font.bold = True
        cl_frame.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)
        
        # 値
        cv_box = slide6.shapes.add_textbox(Inches(6), cy, Inches(4), Inches(0.5))
        cv_frame = cv_box.text_frame
        cv_frame.text = value
        cv_frame.paragraphs[0].font.size = Pt(18)
        cv_frame.paragraphs[0].font.color.rgb = vcolor
        
        cy += Inches(0.65)
    
    # フッターメッセージ
    footer = slide6.shapes.add_textbox(
        Inches(0), Inches(6.8), Inches(13.333), Inches(0.5)
    )
    footer_frame = footer.text_frame
    footer_frame.text = "ご連絡お待ちしております"
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer_frame.paragraphs[0].font.italic = True
    
    # 保存
    output_path = '/root/.openclaw/workspace/language_rpg/株式会社エイジレス_会社紹介_プレミアム.pptx'
    prs.save(output_path)
    print(f"プレミアムプレゼンテーションを作成しました: {output_path}")
    print("\n特徴:")
    print("- グラデーション背景（擬似）")
    print("- シャドウ付きテキスト")
    print("- 3D風カードデザイン")
    print("- アクセントカラー統一")
    print("- 透明度効果")
    print("- レスポンシブな配置")

if __name__ == "__main__":
    create_premium_presentation()
