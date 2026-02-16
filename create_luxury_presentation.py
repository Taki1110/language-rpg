from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from lxml import etree
import io

def create_premium_presentation_v2():
    """プレミアム品質のプレゼンテーションを作成 - 全面刷新版"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 高級感のあるカラーパレット
    DARK_NAVY = RGBColor(15, 23, 42)      # 深い紺色
    GOLD = RGBColor(212, 175, 55)         # ゴールド
    CREAM = RGBColor(254, 250, 240)       # クリーム色
    SLATE = RGBColor(71, 85, 105)         # スレートグレー
    LIGHT_GOLD = RGBColor(251, 240, 213)  # 薄いゴールド
    WHITE = RGBColor(255, 255, 255)
    
    # ========== スライド1: タイトル（映画ポスター風） ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景 - 深い紺色
    bg1 = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_NAVY
    bg1.line.fill.background()
    
    # 装飾的なゴールドライン（上部）
    top_line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.3), Inches(9.333), Pt(2)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = GOLD
    top_line.line.fill.background()
    
    # 会社名 - 大きく、中央に
    title1 = slide1.shapes.add_textbox(Inches(0), Inches(2.6), Inches(13.333), Inches(1.2))
    tf1 = title1.text_frame
    tf1.text = "株式会社エイジレス"
    tf1.paragraphs[0].font.size = Pt(72)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = WHITE
    tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf1.paragraphs[0].font.name = "Hiragino Sans"
    
    # 英語名
    subtitle1 = slide1.shapes.add_textbox(Inches(0), Inches(3.9), Inches(13.333), Inches(0.6))
    stf1 = subtitle1.text_frame
    stf1.text = "A G E L E S S   I N C ."
    stf1.paragraphs[0].font.size = Pt(24)
    stf1.paragraphs[0].font.color.rgb = GOLD
    stf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    stf1.paragraphs[0].font.letter_spacing = Pt(8)
    
    # 装飾的なゴールドライン（下部）
    bottom_line = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.6), Inches(9.333), Pt(2)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = GOLD
    bottom_line.line.fill.background()
    
    # キャッチコピー
    catch1 = slide1.shapes.add_textbox(Inches(0), Inches(5.2), Inches(13.333), Inches(0.5))
    cf1 = catch1.text_frame
    cf1.text = "年齢を超えた価値創造"
    cf1.paragraphs[0].font.size = Pt(20)
    cf1.paragraphs[0].font.color.rgb = LIGHT_GOLD
    cf1.paragraphs[0].alignment = PP_ALIGN.CENTER
    cf1.paragraphs[0].font.italic = True
    
    # ========== スライド2: 会社概要（左右分割レイアウト） ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 左側 - 紺色背景
    left_bg = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5), Inches(7.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = DARK_NAVY
    left_bg.line.fill.background()
    
    # 左側タイトル
    left_title = slide2.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(1))
    ltf = left_title.text_frame
    ltf.text = "会社概要"
    ltf.paragraphs[0].font.size = Pt(48)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = WHITE
    ltf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # 左側英語
    left_en = slide2.shapes.add_textbox(Inches(0.5), Inches(4), Inches(4), Inches(0.5))
    letf = left_en.text_frame
    letf.text = "COMPANY PROFILE"
    letf.paragraphs[0].font.size = Pt(14)
    letf.paragraphs[0].font.color.rgb = GOLD
    letf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # 右側 - 白背景
    right_bg = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(8.333), Inches(7.5)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = CREAM
    right_bg.line.fill.background()
    
    # ミッションステートメント
    mission_box = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(0.8), Inches(7.333), Inches(2.2)
    )
    mission_box.fill.solid()
    mission_box.fill.fore_color.rgb = WHITE
    mission_box.line.color.rgb = GOLD
    mission_box.line.width = Pt(1)
    
    mission_title = slide2.shapes.add_textbox(Inches(5.8), Inches(1), Inches(6.5), Inches(0.5))
    mttf = mission_title.text_frame
    mttf.text = "ミッション"
    mttf.paragraphs[0].font.size = Pt(18)
    mttf.paragraphs[0].font.bold = True
    mttf.paragraphs[0].font.color.rgb = DARK_NAVY
    
    mission_text = slide2.shapes.add_textbox(Inches(5.8), Inches(1.5), Inches(6.5), Inches(1.3))
    mtf = mission_text.text_frame
    mtf.word_wrap = True
    mtf.text = "年齢にとらわれない価値の創造を使命とし、テクノロジーと人間性の融合により、すべての世代が活躍できる社会の実現を目指す。"
    mtf.paragraphs[0].font.size = Pt(14)
    mtf.paragraphs[0].font.color.rgb = SLATE
    mtf.paragraphs[0].line_spacing = 1.3
    
    # 会社情報
    info_items = [
        ("設立", "20XX年"),
        ("従業員数", "XX名"),
        ("資本金", "X,XXX万円"),
        ("本社所在地", "東京都XXX区XXX")
    ]
    
    y_pos = Inches(3.3)
    for label, value in info_items:
        # ラベル
        label_box = slide2.shapes.add_textbox(Inches(5.8), y_pos, Inches(2), Inches(0.4))
        lbf = label_box.text_frame
        lbf.text = label
        lbf.paragraphs[0].font.size = Pt(12)
        lbf.paragraphs[0].font.color.rgb = SLATE
        
        # 値
        value_box = slide2.shapes.add_textbox(Inches(8), y_pos, Inches(4.5), Inches(0.4))
        vbf = value_box.text_frame
        vbf.text = value
        vbf.paragraphs[0].font.size = Pt(14)
        vbf.paragraphs[0].font.bold = True
        vbf.paragraphs[0].font.color.rgb = DARK_NAVY
        
        # セパレータライン
        sep = slide2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.8), y_pos + Inches(0.45), Inches(6.5), Pt(1)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(220, 220, 220)
        sep.line.fill.background()
        
        y_pos += Inches(0.8)
    
    # ========== スライド3: ビジョン（中央配置・インパクト重視） ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = CREAM
    bg3.line.fill.background()
    
    # ヘッダー
    header3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header3.fill.solid()
    header3.fill.fore_color.rgb = DARK_NAVY
    header3.line.fill.background()
    
    header_text = slide3.shapes.add_textbox(Inches(0), Inches(0.35), Inches(13.333), Inches(0.6))
    htf3 = header_text.text_frame
    htf3.text = "ビジョン"
    htf3.paragraphs[0].font.size = Pt(32)
    htf3.paragraphs[0].font.bold = True
    htf3.paragraphs[0].font.color.rgb = WHITE
    htf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 中央のビジョンカード
    vision_card = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(9.333), Inches(3.5)
    )
    vision_card.fill.solid()
    vision_card.fill.fore_color.rgb = WHITE
    vision_card.line.color.rgb = GOLD
    vision_card.line.width = Pt(2)
    
    # ゴールドアクセント（上部）
    v_accent = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(9.333), Pt(6)
    )
    v_accent.fill.solid()
    v_accent.fill.fore_color.rgb = GOLD
    v_accent.line.fill.background()
    
    # ビジョンテキスト
    vision_lines = [
        ("年齢という境界を超え、", Pt(24), False, SLATE),
        ("すべての人が可能性を", Pt(24), False, SLATE),
        ("最大限に発揮できる", Pt(24), False, SLATE),
        ("世界の創造", Pt(36), True, DARK_NAVY)
    ]
    
    y_v = Inches(2.5)
    for text, size, bold, color in vision_lines:
        v_box = slide3.shapes.add_textbox(Inches(2.5), y_v, Inches(8.333), Inches(0.6))
        vtf = v_box.text_frame
        vtf.text = text
        vtf.paragraphs[0].font.size = size
        vtf.paragraphs[0].font.bold = bold
        vtf.paragraphs[0].font.color.rgb = color
        vtf.paragraphs[0].alignment = PP_ALIGN.CENTER
        y_v += Inches(0.6)
    
    # キーワード
    keywords = ["#多様性", "#イノベーション", "#持続可能性"]
    kw_x = Inches(3.5)
    for kw in keywords:
        kw_box = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, kw_x, Inches(6), Inches(1.8), Inches(0.5)
        )
        kw_box.fill.solid()
        kw_box.fill.fore_color.rgb = LIGHT_GOLD
        kw_box.line.fill.background()
        
        kw_text = slide3.shapes.add_textbox(kw_x, Inches(6.05), Inches(1.8), Inches(0.4))
        kw_tf = kw_text.text_frame
        kw_tf.text = kw
        kw_tf.paragraphs[0].font.size = Pt(12)
        kw_tf.paragraphs[0].font.color.rgb = DARK_NAVY
        kw_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        kw_x += Inches(2.2)
    
    # ========== スライド4: 事業内容（3カード横並び） ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = WHITE
    bg4.line.fill.background()
    
    # ヘッダー
    header4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3)
    )
    header4.fill.solid()
    header4.fill.fore_color.rgb = DARK_NAVY
    header4.line.fill.background()
    
    h4_text = slide4.shapes.add_textbox(Inches(0), Inches(0.4), Inches(13.333), Inches(0.6))
    h4tf = h4_text.text_frame
    h4tf.text = "事業内容"
    h4tf.paragraphs[0].font.size = Pt(36)
    h4tf.paragraphs[0].font.bold = True
    h4tf.paragraphs[0].font.color.rgb = WHITE
    h4tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # サービスカード
    services = [
        {
            "icon": "01",
            "title": "デジタル変革",
            "desc": "AI・クラウド・データ分析など最新技術を活用し、企業のビジネスプロセスを根本から変革します。",
            "color": DARK_NAVY
        },
        {
            "icon": "02",
            "title": "人材開発",
            "desc": "多世代が学び合う共生型の人材育成プログラムで、組織の持続的成長を支援します。",
            "color": GOLD
        },
        {
            "icon": "03",
            "title": "グローバル展開",
            "desc": "言語と文化の壁を超えたグローバルビジネス支援で、世界市場での成功を目指します。",
            "color": SLATE
        }
    ]
    
    card_w = Inches(3.8)
    start_x = Inches(0.8)
    gap = Inches(0.4)
    
    for i, svc in enumerate(services):
        x = start_x + i * (card_w + gap)
        
        # カード背景
        card = slide4.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), card_w, Inches(5.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM
        card.line.fill.background()
        
        # 番号サークル
        circle = slide4.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.4), Inches(2.2), Inches(1), Inches(1)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = svc["color"]
        circle.line.fill.background()
        
        num = slide4.shapes.add_textbox(x + Inches(1.4), Inches(2.4), Inches(1), Inches(0.6))
        ntf = num.text_frame
        ntf.text = svc["icon"]
        ntf.paragraphs[0].font.size = Pt(28)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = WHITE
        ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # タイトル
        title = slide4.shapes.add_textbox(x, Inches(3.5), card_w, Inches(0.6))
        ttf = title.text_frame
        ttf.text = svc["title"]
        ttf.paragraphs[0].font.size = Pt(24)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = svc["color"]
        ttf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 説明
        desc = slide4.shapes.add_textbox(x + Inches(0.3), Inches(4.2), card_w - Inches(0.6), Inches(2.5))
        dtf = desc.text_frame
        dtf.word_wrap = True
        dtf.text = svc["desc"]
        dtf.paragraphs[0].font.size = Pt(13)
        dtf.paragraphs[0].font.color.rgb = SLATE
        dtf.paragraphs[0].alignment = PP_ALIGN.CENTER
        dtf.paragraphs[0].line_spacing = 1.4
    
    # ========== スライド5: 強み（4つのポイント） ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = DARK_NAVY
    bg5.line.fill.background()
    
    # タイトル
    title5 = slide5.shapes.add_textbox(Inches(0), Inches(0.5), Inches(13.333), Inches(0.8))
    t5f = title5.text_frame
    t5f.text = "私たちの強み"
    t5f.paragraphs[0].font.size = Pt(40)
    t5f.paragraphs[0].font.bold = True
    t5f.paragraphs[0].font.color.rgb = WHITE
    t5f.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    sub5 = slide5.shapes.add_textbox(Inches(0), Inches(1.2), Inches(13.333), Inches(0.4))
    s5f = sub5.text_frame
    s5f.text = "OUR STRENGTHS"
    s5f.paragraphs[0].font.size = Pt(14)
    s5f.paragraphs[0].font.color.rgb = GOLD
    s5f.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 強みポイント
    strengths = [
        ("01", "多様な世代の知見を\n活かしたイノベーション", "20年以上の実務経験を持つベテランと、最新技術に精通した若手が協働"),
        ("02", "最新テクノロジーと\n実務経験の融合", "理論と実践のバランスを取れたソリューションを提供"),
        ("03", "お客様に寄り添う\n柔軟な対応力", "大手にはない機動力で、お客様のニーズに迅速に対応"),
        ("04", "持続可能な社会への\n貢献", "SDGsの達成に向けた取り組みを事業の核に据える")
    ]
    
    positions = [
        (Inches(1), Inches(2)),
        (Inches(7), Inches(2)),
        (Inches(1), Inches(4.5)),
        (Inches(7), Inches(4.5))
    ]
    
    for i, (num, title, desc) in enumerate(strengths):
        x, y = positions[i]
        
        # カード背景
        card = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.3), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59)
        card.line.color.rgb = GOLD
        card.line.width = Pt(1)
        
        # 番号
        num_box = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.6))
        ntf = num_box.text_frame
        ntf.text = num
        ntf.paragraphs[0].font.size = Pt(32)
        ntf.paragraphs[0].font.bold = True
        ntf.paragraphs[0].font.color.rgb = GOLD
        
        # タイトル
        title_box = slide5.shapes.add_textbox(x + Inches(1.2), y + Inches(0.3), Inches(3.8), Inches(0.8))
        ttf = title_box.text_frame
        ttf.word_wrap = True
        ttf.text = title
        ttf.paragraphs[0].font.size = Pt(16)
        ttf.paragraphs[0].font.bold = True
        ttf.paragraphs[0].font.color.rgb = WHITE
        
        # 説明
        desc_box = slide5.shapes.add_textbox(x + Inches(0.3), y + Inches(1.2), Inches(4.7), Inches(0.8))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dtf.text = desc
        dtf.paragraphs[0].font.size = Pt(11)
        dtf.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
    
    # ========== スライド6: お問い合わせ（シンプル・高級感） ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg6 = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = CREAM
    bg6.line.fill.background()
    
    # 左側 - 紺色パネル
    left_panel = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = DARK_NAVY
    left_panel.line.fill.background()
    
    # 左側タイトル
    left_title = slide6.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(4.5), Inches(1))
    ltf6 = left_title.text_frame
    ltf6.text = "お問い合わせ"
    ltf6.paragraphs[0].font.size = Pt(40)
    ltf6.paragraphs[0].font.bold = True
    ltf6.paragraphs[0].font.color.rgb = WHITE
    
    left_sub = slide6.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(4.5), Inches(0.5))
    lsf6 = left_sub.text_frame
    lsf6.text = "CONTACT US"
    lsf6.paragraphs[0].font.size = Pt(14)
    lsf6.paragraphs[0].font.color.rgb = GOLD
    
    # 右側 - 連絡先情報
    contact_info = [
        ("会社名", "株式会社エイジレス"),
        ("所在地", "〒XXX-XXXX 東京都XXX区XXX"),
        ("TEL", "03-XXXX-XXXX"),
        ("Email", "info@ageless.co.jp"),
        ("Web", "www.ageless.co.jp")
    ]
    
    y_c = Inches(1.5)
    for label, value in contact_info:
        # ラベル
        cl = slide6.shapes.add_textbox(Inches(6.5), y_c, Inches(2), Inches(0.4))
        clf = cl.text_frame
        clf.text = label
        clf.paragraphs[0].font.size = Pt(12)
        clf.paragraphs[0].font.color.rgb = SLATE
        
        # 値
        cv = slide6.shapes.add_textbox(Inches(6.5), y_c + Inches(0.35), Inches(6), Inches(0.5))
        cvf = cv.text_frame
        cvf.text = value
        cvf.paragraphs[0].font.size = Pt(16)
        cvf.paragraphs[0].font.bold = True
        cvf.paragraphs[0].font.color.rgb = DARK_NAVY
        
        y_c += Inches(1)
    
    # フッターメッセージ
    footer = slide6.shapes.add_textbox(Inches(6.5), Inches(6.5), Inches(6.5), Inches(0.5))
    ff = footer.text_frame
    ff.text = "ご連絡お待ちしております"
    ff.paragraphs[0].font.size = Pt(14)
    ff.paragraphs[0].font.color.rgb = SLATE
    ff.paragraphs[0].font.italic = True
    
    # 保存
    output_path = '/root/.openclaw/workspace/language_rpg/株式会社エイジレス_会社紹介_高級版.pptx'
    prs.save(output_path)
    print(f"高級版プレゼンテーションを作成しました: {output_path}")
    print("\nデザイン特徴:")
    print("- 深い紺色とゴールドの高級感あるカラースキーム")
    print("- 映画ポスター風のタイトルスライド")
    print("- 左右分割レイアウトの会社概要")
    print("- 中央配置のインパクト重視ビジョン")
    print("- 3カード横並びの事業内容")
    print("- 2x2グリッドの強み紹介")
    print("- シンプルかつ高級感あるお問い合わせ")

if __name__ == "__main__":
    create_premium_presentation_v2()
